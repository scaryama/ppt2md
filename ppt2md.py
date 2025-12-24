"""
Script to convert PPT/PPTX files to Markdown format
Runs safely locally and converts all slides of PowerPoint files to Markdown.

Required packages:
    pip install python-pptx
"""

import os
import sys
from pathlib import Path
from pptx import Presentation


def ppt_to_markdown(ppt_path, output_dir=None, log_callback=None):
    """
    Converts all slides of a PPT/PPTX file into a single Markdown file.
    
    Args:
        ppt_path: Path to the PPT/PPTX file to convert
        output_dir: Output directory path (None means same location as original file)
        log_callback: Callback function to output log messages (uses print if None)
    
    Returns:
        bool: Conversion success status
    """
    # Use default print if no callback provided
    log = log_callback if log_callback else print
    
    try:
        # Check if file exists
        if not os.path.exists(ppt_path):
            log(f"❌ File not found: {ppt_path}")
            return False
        
        # Set output directory
        if output_dir is None:
            output_dir = Path(ppt_path).parent
        else:
            output_dir = Path(output_dir)
            output_dir.mkdir(parents=True, exist_ok=True)
        
        # Read PowerPoint file
        log(f"📄 Reading file: {ppt_path}")
        prs = Presentation(ppt_path)
        slide_count = len(prs.slides)
        
        log(f"📊 Number of slides: {slide_count}")
        log(f"📁 Output directory: {output_dir}\n")
        
        # Generate output filename
        base_name = Path(ppt_path).stem
        output_filename = f"{base_name}.md"
        output_path = output_dir / output_filename
        
        # Generate Markdown content
        markdown_content = []
        markdown_content.append(f"# {base_name}\n")
        markdown_content.append(f"*Source file: {Path(ppt_path).name}*\n")
        markdown_content.append(f"*Total slides: {slide_count}*\n")
        markdown_content.append("---\n\n")
        
        # Convert each slide
        for idx, slide in enumerate(prs.slides, 1):
            log(f"  📋 Processing slide {idx}/{slide_count}")
            
            try:
                # Slide header
                markdown_content.append(f"## Slide {idx}\n\n")
                
                # Process slide shapes
                slide_content = process_slide(slide, idx)
                if slide_content:
                    markdown_content.append(slide_content)
                    markdown_content.append("\n")
                else:
                    markdown_content.append("*This slide is empty.*\n\n")
                
                # Add separator between slides
                if idx < slide_count:
                    markdown_content.append("---\n\n")
                
            except Exception as e:
                log(f"    ⚠️  Error processing slide {idx}: {e}")
                import traceback
                traceback_str = traceback.format_exc()
                log(traceback_str)
                markdown_content.append(f"*Error processing this slide: {e}*\n\n")
        
        # Save Markdown file
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(''.join(markdown_content))
        
        log(f"\n✅ Conversion complete: {output_filename}")
        return True
        
    except Exception as e:
        log(f"❌ Error occurred: {e}")
        import traceback
        traceback_str = traceback.format_exc()
        log(traceback_str)
        return False


def process_slide(slide, slide_num):
    """
    Processes a slide and converts its content to Markdown.
    
    Args:
        slide: pptx.slide.Slide object
        slide_num: Slide number for image placeholders
    
    Returns:
        str: Markdown content for the slide
    """
    content_parts = []
    
    for shape in slide.shapes:
        # Process text boxes
        if hasattr(shape, "text") and shape.text.strip():
            text = shape.text.strip()
            # Check if it's a title (usually first shape or has specific formatting)
            if shape == slide.shapes[0] and len(text) < 100:
                content_parts.append(f"### {text}\n\n")
            else:
                # Process paragraphs with bullet points
                paragraphs = text.split('\n')
                for para in paragraphs:
                    para = para.strip()
                    if para:
                        # Check for bullet points (lines starting with •, -, *, etc.)
                        if para.startswith(('•', '-', '*', '◦')):
                            content_parts.append(f"- {para[1:].strip()}\n")
                        else:
                            content_parts.append(f"{para}\n")
                content_parts.append("\n")
        
        # Process tables
        elif hasattr(shape, "table"):
            table_md = table_to_markdown(shape.table)
            if table_md:
                content_parts.append(table_md)
                content_parts.append("\n")
        
        # Process images (placeholder)
        elif hasattr(shape, "image"):
            content_parts.append(f"[Image: Slide {slide_num}]\n\n")
    
    return ''.join(content_parts)


def table_to_markdown(table):
    """
    Converts a PowerPoint table to Markdown table format.
    
    Args:
        table: pptx.table.Table object
    
    Returns:
        str: Markdown table string
    """
    if not table.rows:
        return ""
    
    lines = []
    
    # Process header row (first row)
    header_row = table.rows[0]
    headers = []
    for cell in header_row.cells:
        cell_text = cell.text.strip().replace("|", "\\|").replace("\n", "<br>")
        headers.append(cell_text)
    
    lines.append("| " + " | ".join(headers) + " |")
    
    # Generate separator line
    separator = "| " + " | ".join(["---"] * len(headers)) + " |"
    lines.append(separator)
    
    # Process data rows
    for row in table.rows[1:]:
        cells = []
        for cell in row.cells:
            cell_text = cell.text.strip().replace("|", "\\|").replace("\n", "<br>")
            cells.append(cell_text)
        lines.append("| " + " | ".join(cells) + " |")
    
    return "\n".join(lines) + "\n"


def main():
    """Main function"""
    if len(sys.argv) <= 1:
        return
    ppt_path = sys.argv[1]
    
    # Set output directory (optional)
    output_dir = None
    if len(sys.argv) > 2:
        output_dir = sys.argv[2]
    
    print("=" * 60)
    print("PPT to Markdown Converter")
    print("=" * 60)
    
    success = ppt_to_markdown(ppt_path, output_dir, log_callback=None)
    
    if success:
        print("\n✨ Conversion completed successfully!")
    else:
        print("\n❌ Conversion failed.")
        sys.exit(1)


if __name__ == "__main__":
    main()

